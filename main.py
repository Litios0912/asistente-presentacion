import os
import json
import uuid
import shutil
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, WebSocket, WebSocketDisconnect
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse
from dotenv import load_dotenv

load_dotenv()

app = FastAPI(title="Asistente de Presentación")

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

sessions = {}

class PresentationSession:
    def __init__(self):
        self.slides = []
        self.current_slide = 0
        self.title = ""
        self.transcript_history = []
        self.groq_key = None
        self.analysis = None
        self.search_cache = {}

    def extract_slides(self, filepath: str):
        ext = Path(filepath).suffix.lower()
        if ext == ".pdf":
            self._extract_pdf(filepath)
        elif ext == ".pptx":
            self._extract_pptx(filepath)

    def _extract_pdf(self, filepath: str):
        import fitz
        doc = fitz.open(filepath)
        for page in doc:
            text = page.get_text().strip()
            if text:
                self.slides.append({"number": len(self.slides) + 1, "content": text})
        doc.close()

    def _extract_pptx(self, filepath: str):
        from pptx import Presentation
        prs = Presentation(filepath)
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
            text = "\n".join(parts)
            if text:
                self.slides.append({"number": len(self.slides) + 1, "content": text})


async def web_search(query: str, max_results: int = 5) -> list:
    try:
        from duckduckgo_search import DDGS
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                results.append(r.get("body", ""))
        return results[:max_results]
    except Exception:
        return []


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    session_id = str(uuid.uuid4())
    ext = Path(file.filename).suffix.lower()
    if ext not in (".pdf", ".pptx"):
        return JSONResponse({"error": "Solo PDF y PPTX son soportados"}, status_code=400)

    dest = UPLOAD_DIR / f"{session_id}{ext}"
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)

    session = PresentationSession()
    session.title = Path(file.filename).stem
    session.extract_slides(str(dest))

    if not session.slides:
        return JSONResponse({"error": "No se pudo extraer texto de la presentación"}, status_code=400)

    sessions[session_id] = session
    return {
        "session_id": session_id,
        "title": session.title,
        "total_slides": len(session.slides),
        "slides": session.slides,
    }


@app.post("/analyze/{session_id}")
async def analyze_presentation(session_id: str, data: dict):
    session = sessions.get(session_id)
    if not session:
        return JSONResponse({"error": "Sesión no encontrada"}, status_code=404)

    groq_key = data.get("groq_key") or session.groq_key or GROQ_API_KEY
    if not groq_key:
        return JSONResponse({"error": "GROQ_API_KEY no configurada"}, status_code=400)

    from groq import AsyncGroq
    client = AsyncGroq(api_key=groq_key)
    session.groq_key = groq_key

    slides_text = "\n\n".join(
        f"--- Slide {s['number']} ---\n{s['content'][:1000]}"
        for s in session.slides
    )

    try:
        completion = await client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": """Dada una presentación, analiza cada slide y extrae:

1. Temas principales del slide (para búsqueda web)
2. Puntos clave que vale la pena mencionar
3. Preguntas que la audiencia podría hacer
4. Respuestas sugeridas para esas preguntas

Responde solo JSON:
{
  "slides": [
    {
      "number": 1,
      "key_points": ["punto clave 1", "punto clave 2"],
      "search_terms": ["términos para buscar en internet"],
      "questions": [
        {"q": "pregunta", "a": "respuesta"}
      ]
    }
  ]
}"""},
                {"role": "user", "content": f"Título: {session.title}\n\n{slides_text}"},
            ],
            temperature=0.2,
            max_tokens=2000,
        )

        raw = completion.choices[0].message.content.strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        session.analysis = json.loads(raw)

        search_terms = set()
        for s in session.analysis.get("slides", []):
            for t in s.get("search_terms", []):
                search_terms.add(t)

        for term in list(search_terms)[:3]:
            try:
                results = await web_search(term)
                for s in session.analysis.get("slides", []):
                    if term in s.get("search_terms", []):
                        s.setdefault("web_results", []).extend(results[:3])
            except Exception:
                pass

        return {"status": "ok", "analysis": session.analysis}

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.websocket("/ws/{session_id}")
async def websocket_endpoint(websocket: WebSocket, session_id: str):
    await websocket.accept()
    session = sessions.get(session_id)
    if not session:
        await websocket.send_json({"error": "Sesión no encontrada"})
        await websocket.close()
        return

    client = None

    try:
        while True:
            data = await websocket.receive_json()
            action = data.get("action")

            if action == "set_key":
                key = data.get("key", "").strip()
                if key:
                    session.groq_key = key
                    await websocket.send_json({"type": "key_set"})

            elif action == "transcribe":
                if not client:
                    api_key = session.groq_key or GROQ_API_KEY
                    if not api_key:
                        await websocket.send_json({"type": "error", "message": "Configura tu API key de Groq primero"})
                        continue
                    from groq import AsyncGroq
                    client = AsyncGroq(api_key=api_key)

                transcript = data.get("text", "").strip()
                if not transcript:
                    continue

                session.transcript_history.append(transcript)
                session.transcript_history = session.transcript_history[-20:]

                current = session.current_slide
                slide = session.slides[current] if current < len(session.slides) else None
                slide_content = slide["content"] if slide else ""

                recent = " ".join(session.transcript_history[-5:])

                is_question = bool(transcript.strip().endswith("?"))

                slide_analysis = None
                web_results = []
                search_terms = []
                qa_list = []
                if session.analysis:
                    for sa in session.analysis.get("slides", []):
                        if sa["number"] == current + 1:
                            slide_analysis = sa
                            search_terms = sa.get("search_terms", [])
                            qa_list = sa.get("questions", [])
                            web_results = sa.get("web_results", [])
                            break

                if is_question:
                    system_prompt = """El presentador recibió una pregunta. Tu función es dar información útil para responderla bien.

Solo responde JSON:
{"advice":"información para responder la pregunta (máx 3 oraciones)","slide_match":true,"off_topic":false,"confidence":0.9,"extra_info":["dato relevante 1","dato relevante 2"],"suggested_questions":[],"suggested_slide":null,"is_question":true,"qa_answer":"respuesta basada en el contenido"}

- advice: dato útil para responder
- qa_answer: mejor respuesta posible según el contenido
- extra_info: datos adicionales relevantes a la pregunta"""

                    user_msg = json.dumps({
                        "current_slide": current + 1,
                        "slide_content": slide_content[:1500],
                        "pregunta": transcript,
                        "presentation_title": session.title,
                        "qa_preparados": qa_list,
                        "resultados_web": web_results,
                    })
                else:
                    system_prompt = """Eres un asistente que ENRIQUECE presentaciones. Tu función es AÑADIR información relevante a lo que el presentador dice.

NO evalúes si lo que dice coincide con el slide. En lugar de eso, PROPORCIONA:
- Contexto adicional, datos interesantes o conexiones con otros temas
- Información extra que la audiencia encontraría valiosa
- Posibles preguntas que podrían surgir

Solo responde JSON:
{"advice":"información enriquecida sobre el tema que está tratando (máx 3 oraciones)","slide_match":true,"off_topic":false,"confidence":0.9,"extra_info":["información adicional relevante 1","información 2","información 3"],"suggested_questions":["pregunta que podría hacer la audiencia 1?","pregunta 2?"],"key_points_covered":[],"key_points_missed":[],"suggested_slide":null,"is_question":false,"qa_answer":null}

- advice: INFORMACIÓN ENRIQUECIDA sobre el tema que está presentando. Basada en el contenido del slide, lo que dijo, los resultados de búsqueda web y los Q&A preparados.
- extra_info: lista de 2-4 datos adicionales, contextos o conexiones interesantes
- suggested_questions: 1-2 preguntas relevantes que la audiencia podría hacer
- slide_match: siempre true
- off_topic: siempre false
- confidence: siempre 0.9"""

                    all_summaries = [
                        f"Slide {s['number']}: {s['content'][:200]}"
                        for s in session.slides
                    ]

                    user_msg = json.dumps({
                        "current_slide": current + 1,
                        "total_slides": len(session.slides),
                        "slide_content": slide_content[:1500],
                        "lo_que_dijo": recent,
                        "presentation_title": session.title,
                        "all_slides_summary": all_summaries,
                        "resultados_busqueda_web": web_results,
                        "puntos_clave": slide_analysis["key_points"] if slide_analysis else [],
                        "qa_preparados": qa_list,
                    })

                try:
                    completion = await client.chat.completions.create(
                        model="llama-3.1-8b-instant",
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_msg},
                        ],
                        temperature=0.4,
                        max_tokens=500,
                    )

                    raw = completion.choices[0].message.content.strip()
                    raw = raw.replace("```json", "").replace("```", "").strip()
                    result = json.loads(raw)

                    if result.get("suggested_slide") and isinstance(result["suggested_slide"], (int, float)):
                        suggested = int(result["suggested_slide"])
                        if 1 <= suggested <= len(session.slides):
                            session.current_slide = suggested - 1
                            result["current_slide"] = suggested

                    result["current_slide"] = session.current_slide + 1
                    result.setdefault("extra_info", [])
                    result.setdefault("suggested_questions", [])
                    result.setdefault("key_points_covered", [])
                    result.setdefault("key_points_missed", [])

                    await websocket.send_json({"type": "advice", **result})

                except Exception as e:
                    await websocket.send_json({"type": "error", "message": str(e)})

            elif action == "search":
                query = data.get("query", "").strip()
                if query:
                    results = await web_search(query)
                    await websocket.send_json({"type": "search_results", "query": query, "results": results})

            elif action == "set_slide":
                slide_num = data.get("slide", 1) - 1
                if 0 <= slide_num < len(session.slides):
                    session.current_slide = slide_num
                    await websocket.send_json({
                        "type": "slide_changed",
                        "current_slide": slide_num + 1,
                    })

    except WebSocketDisconnect:
        pass
    except Exception as e:
        try:
            await websocket.send_json({"type": "error", "message": str(e)})
        except:
            pass


if not GROQ_API_KEY:
    print("⚠️  ADVERTENCIA: GROQ_API_KEY no está configurada. Usa el archivo .env")

app.mount("/", StaticFiles(directory="static", html=True), name="static")
